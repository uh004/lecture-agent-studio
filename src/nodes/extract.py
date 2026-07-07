import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from src.core.state import State
from src.core.utils import clean_text, export_slide_as_png


def node_parse_all(state: State) -> State:
    """v3 PPT 파싱: 제목, 본문, 노트, 표, 이미지, 도형 텍스트, 링크, PNG 스냅샷 추출"""
    pptx_path = state["pptx_path"]
    work_dir = state.get("work_dir", "./")

    base_dir = Path(work_dir).expanduser().resolve()
    slides_dir = base_dir / "slides"
    media_dir = base_dir / "media"
    slides_dir.mkdir(parents=True, exist_ok=True)
    media_dir.mkdir(parents=True, exist_ok=True)

    if not os.path.exists(pptx_path):
        print(f"❌ 오류: '{pptx_path}' 파일이 존재하지 않습니다.")
        return state

    prs = Presentation(pptx_path)
    slide_count = len(prs.slides)
    print(f"총 {slide_count}장의 슬라이드를 처리합니다.")

    titles_list, texts_list, notes_list = [], [], []
    tables_list, images_list, snapshots_list = [], [], []
    shape_texts_list, links_list = [], []
    url_pattern = r"https?://[^\s]+"

    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n=== {slide_idx+1}번째 슬라이드 처리 중 ===")

        # 1) PNG 스냅샷 생성
        slide_state = {"pptx_path": pptx_path, "work_dir": str(slides_dir), "slide_index": slide_idx}
        snapshot_path = ""
        try:
            slide_state = export_slide_as_png(slide_state)
            src_png = slide_state.get("slide_image", "")
            dst_png = slides_dir / f"slide_{slide_idx+1}.png"
            if src_png and os.path.exists(src_png):
                os.replace(src_png, dst_png)
                snapshot_path = str(dst_png)
        except Exception as e:
            print(f"  ⚠️ PNG 변환 실패 (텍스트 추출은 계속): {e}")

        # v3: 발표자 노트 추출
        slide_note = ""
        try:
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                if notes_tf and notes_tf.text.strip():
                    slide_note = notes_tf.text.strip()
        except Exception:
            pass

        slide_title, body_texts, shape_texts = "", [], []
        slide_tables, slide_images = [], []
        slide_links = set()

        # 차트 데이터 추출
        for sh in slide.shapes:
            if sh.has_chart:
                try:
                    chart = sh.chart
                    chart_data = ["\n[📊 차트 데이터]"]
                    if chart.has_title and chart.chart_title.has_text_frame:
                        chart_data.append(f"차트 제목: {chart.chart_title.text_frame.text}")
                    series_names = [s.name for s in chart.series]
                    chart_data.append("항목 | " + " | ".join(series_names))
                    if chart.plots:
                        plot = chart.plots[0]
                        categories = [c.label for c in plot.categories]
                        for i, cat in enumerate(categories):
                            row_vals = [str(cat)]
                            for s in chart.series:
                                val = s.values[i] if i < len(s.values) else ""
                                row_vals.append(str(val))
                            chart_data.append(" | ".join(row_vals))
                    body_texts.append("\n".join(chart_data))
                except Exception as e:
                    print(f"  ⚠️ 차트 데이터 추출 실패: {e}")

        # 제목 추출
        for sh in slide.shapes:
            try:
                if sh.is_placeholder and sh.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                    if sh.has_text_frame and sh.text.strip():
                        slide_title = sh.text.strip()
            except:
                pass

        # 도형 텍스트 (그룹 포함 재귀)
        def collect_shape_texts(shape):
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub in shape.shapes:
                    collect_shape_texts(sub)
                return
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    shape_texts.append(t)
        for sh in slide.shapes:
            collect_shape_texts(sh)

        # 링크 추출
        for sh in slide.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    runs_text = "".join(r.text for r in p.runs)
                    found = re.findall(url_pattern, runs_text)
                    if found:
                        slide_links.update(found)
            try:
                if hasattr(sh, "click_action") and sh.click_action.hyperlink.address:
                    slide_links.add(sh.click_action.hyperlink.address)
            except:
                pass

        # 본문 텍스트 (제목/도형/링크 제외)
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            try:
                if sh.is_placeholder and sh.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                    continue
            except:
                pass
            for p in sh.text_frame.paragraphs:
                txt = "".join(r.text for r in p.runs).strip()
                if not txt or txt in shape_texts or re.search(url_pattern, txt):
                    continue
                body_texts.append(txt)

        # 표 추출
        for sh in slide.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.TABLE:
                tbl = [[cell.text.strip() for cell in row.cells] for row in sh.table.rows]
                slide_tables.append(tbl)

        # 이미지 추출 (그룹/숨겨진 이미지 포함)
        def extract_images_from_shape(shape, img_list):
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in shape.shapes:
                    extract_images_from_shape(sub_shape, img_list)
                return
            if hasattr(shape, "image"):
                try:
                    ext = shape.image.ext
                except:
                    ext = "png"
                img_name = f"slide{slide_idx+1}_img_{len(img_list)+1}.{ext}"
                img_path = media_dir / img_name
                with open(img_path, "wb") as f:
                    f.write(shape.image.blob)
                img_list.append(str(img_path))
        for sh in slide.shapes:
            extract_images_from_shape(sh, slide_images)

        # 리스트에 저장
        titles_list.append(slide_title)
        texts_list.append("\n".join(body_texts).strip())
        notes_list.append(slide_note)
        tables_list.append(slide_tables)
        images_list.append(slide_images)
        snapshots_list.append(snapshot_path)
        shape_texts_list.append(list(dict.fromkeys(shape_texts)))
        links_list.append(list(slide_links))

        print(f"  → 제목:{slide_title}, 본문:{len(body_texts)}, 노트:{'있음' if slide_note else '없음'}, 표:{len(slide_tables)}, 이미지:{len(slide_images)}")

    state["total_slides"] = slide_count
    state["titles"] = titles_list
    state["texts"] = texts_list
    state["notes"] = notes_list
    state["tables"] = tables_list
    state["images"] = images_list
    state["slide_image"] = snapshots_list
    state["shape_texts"] = shape_texts_list
    state["links"] = links_list
    return state
